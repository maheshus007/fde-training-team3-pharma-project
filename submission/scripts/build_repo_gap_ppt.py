#!/usr/bin/env python3
"""Stakeholder briefing: inherited package, gaps, close-out approach."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parents[1] / "artefacts" / "AEGIS_PHARMA_STAKEHOLDER_BRIEFING.pptx"

NAVY = RGBColor(0x0A, 0x16, 0x28)
NAVY2 = RGBColor(0x12, 0x2A, 0x45)
TEAL = RGBColor(0x0E, 0x4D, 0x5A)
GOLD = RGBColor(0xB8, 0x95, 0x2C)
PAPER = RGBColor(0xF4, 0xF1, 0xEA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1A, 0x1F, 0x26)
MUTED = RGBColor(0x5A, 0x62, 0x70)
RULE = RGBColor(0xD8, 0xD1, 0xC4)
ROSE = RGBColor(0x7A, 0x2E, 0x2E)
OK = RGBColor(0x1F, 0x5C, 0x45)
SOFT = RGBColor(0xE8, 0xE3, 0xD8)

W = Inches(13.333)
H = Inches(7.5)
TOTAL = 14


def _run(p, text, size, bold=False, color=INK, name="Calibri"):
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    return run


def tf_box(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor])
    except Exception:
        pass
    return tf


def line(tf, text, size=16, bold=False, color=INK, after=6, align=PP_ALIGN.LEFT, first=False):
    p = tf.paragraphs[0] if first or (not tf.paragraphs[0].runs and not tf.paragraphs[0].text) else tf.add_paragraph()
    if first:
        p = tf.paragraphs[0]
        if p.runs:
            p = tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(after)
    p.space_before = Pt(0)
    _run(p, text, size, bold, color)
    return p


def put(slide, l, t, w, h, text, size=16, bold=False, color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tf = tf_box(slide, l, t, w, h, anchor)
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    _run(p, text, size, bold, color)
    return tf


def rect(slide, l, t, w, h, fill, line_color=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line_color:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(0.75)
    else:
        sh.line.fill.background()
    return sh


def hairline(slide, l, t, w, color=RULE):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, Pt(1.1))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


def footer(slide, n: int) -> None:
    rect(slide, Inches(0), Inches(7.22), W, Inches(0.28), NAVY)
    put(slide, Inches(0.55), Inches(7.24), Inches(8.5), Inches(0.24),
        "AEGIS-PHARMA  ·  Stakeholder briefing  ·  Confidential  ·  16 August 2026",
        10, False, RGBColor(0xC9, 0xD0, 0xD8), PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    put(slide, Inches(11.4), Inches(7.24), Inches(1.4), Inches(0.24),
        f"{n:02d}  /  {TOTAL:02d}", 10, False, GOLD, PP_ALIGN.RIGHT, MSO_ANCHOR.MIDDLE)


def content_head(slide, section: str, title: str, deck: str | None = None) -> None:
    rect(slide, Inches(0), Inches(0), W, Inches(1.22), NAVY)
    rect(slide, Inches(0), Inches(0), Inches(0.12), Inches(1.22), GOLD)
    put(slide, Inches(0.55), Inches(0.18), Inches(12.2), Inches(0.28),
        section.upper(), 11, True, GOLD)
    put(slide, Inches(0.55), Inches(0.48), Inches(12.2), Inches(0.58),
        title, 26, True, WHITE)
    if deck:
        put(slide, Inches(0.55), Inches(1.38), Inches(12.2), Inches(0.38),
            deck, 14, False, MUTED)


def blank() -> object:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, Inches(0), Inches(0), W, H, PAPER)
    return s


def style_table(table, header=NAVY, zebra=True):
    table.first_row = True
    for cell in table.rows[0].cells:
        cell.fill.solid()
        cell.fill.fore_color.rgb = header
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(12)
                r.font.color.rgb = WHITE
                r.font.name = "Calibri"
    for i in range(1, len(table.rows)):
        row = table.rows[i]
        bg = WHITE if i % 2 else SOFT
        for cell in row.cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(12)
                    r.font.color.rgb = INK
                    r.font.name = "Calibri"


def set_cell(table, r, c, text, bold=False):
    cell = table.cell(r, c)
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    _run(p, text, 12, bold, WHITE if r == 0 else INK)


def build() -> Path:
    global prs
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # ------------------------------------------------------------------ 01 Title
    s = blank()
    rect(s, Inches(0), Inches(0), W, H, NAVY)
    rect(s, Inches(0), Inches(0), Inches(0.14), H, GOLD)
    put(s, Inches(0.75), Inches(1.35), Inches(11.8), Inches(0.32),
        "PROJECT AEGIS-PHARMA   ·   TEAM 3   ·   STAKEHOLDER BRIEFING", 13, True, GOLD)
    put(s, Inches(0.75), Inches(1.85), Inches(11.8), Inches(1.7),
        "Integrity of the inherited package,\nthe gaps that mattered, and the\nclose-out we recommend.", 36, True, WHITE)
    hairline(s, Inches(0.75), Inches(3.85), Inches(2.2), GOLD)
    put(s, Inches(0.75), Inches(4.1), Inches(11.5), Inches(1.15),
        "An inspectable account of what the monolithic challenge repository already contains,\nwhat was not yet defensible in the participant layer, and how those gaps were closed\nwithout inventing regulated evidence or expanding prohibited write paths.",
        16, False, RGBColor(0xC5, 0xCD, 0xD6))
    put(s, Inches(0.75), Inches(6.35), Inches(8), Inches(0.35),
        "16 August 2026     ·     Offline synthetic package     ·     Advisory POC only",
        13, False, GOLD)
    put(s, Inches(9.4), Inches(6.35), Inches(3.2), Inches(0.35),
        "Decision support  ·  not a go-live", 13, False, RGBColor(0xC5, 0xCD, 0xD6), PP_ALIGN.RIGHT)
    notes(s, "Open by stating this is a decision briefing, not a technical tour. The inherited package is complete as a challenge. Our job was to make Team 3’s layer inspectable: every material inject named, tested, and honest about contradictions. Ask for conditional acceptance of the close-out — not production go.")

    # ------------------------------------------------------------------ 02 Purpose
    s = blank()
    content_head(s, "Purpose of this session", "What we are asking you to take away",
                 "Fifteen minutes. Three conclusions. One residual-risk acceptance.")
    boxes = [
        ("01", "The package is a monolith by design",
         "Challenge evidence and our work share one tree. That is a control, not a limitation. Immutable areas stay untouched."),
        ("02", "The material gaps were inspectability gaps",
         "The catalogue already had 84 injects. What failed was named ownership, executed tests, and an evidence pack a reviewer can read without us in the room."),
        ("03", "We closed them without fabricating evidence",
         "Abstain on research/clinical write paths. Record case/data tensions. Refresh the machine-readable evidence pack. Residual items stay labelled."),
    ]
    for i, (n, title, body) in enumerate(boxes):
        top = Inches(1.95) + Inches(i * 1.6)
        rect(s, Inches(0.55), top, Inches(12.25), Inches(1.45), WHITE, RULE)
        rect(s, Inches(0.55), top, Inches(0.1), Inches(1.45), GOLD if i == 2 else TEAL)
        put(s, Inches(0.85), top + Inches(0.22), Inches(1.0), Inches(0.4), n, 22, True, TEAL)
        put(s, Inches(2.0), top + Inches(0.22), Inches(10.4), Inches(0.4), title, 18, True, NAVY)
        put(s, Inches(2.0), top + Inches(0.68), Inches(10.4), Inches(0.6), body, 14, False, MUTED)
    footer(s, 2)
    notes(s, "Do not walk the file tree yet. Stakeholder first question is ‘what do you want from me?’ Answer: accept the close-out method, accept the residual NI fixtures as labelled, and do not ask us to invent missing CSV rows.")

    # ------------------------------------------------------------------ 03 Agenda
    s = blank()
    content_head(s, "Agenda", "How the briefing is structured")
    rows = [
        ("01", "The inherited package", "What the monolith already is, who may write where, and the hard gates that bound AEGIS."),
        ("02", "The gaps that mattered", "Coverage holes, evidence-pack defects, and challenge contradictions we will not ‘clean’."),
        ("03", "The close-out approach", "The bar we used, what is now closed, what remains open, and the decision requested."),
    ]
    for i, (n, title, body) in enumerate(rows):
        left = Inches(0.55) + Inches(i * 4.2)
        rect(s, left, Inches(1.85), Inches(3.95), Inches(4.7), WHITE, RULE)
        rect(s, left, Inches(1.85), Inches(3.95), Inches(0.08), GOLD)
        put(s, left + Inches(0.28), Inches(2.2), Inches(3.4), Inches(0.45), n, 28, True, TEAL)
        put(s, left + Inches(0.28), Inches(2.8), Inches(3.4), Inches(1.1), title, 20, True, NAVY)
        put(s, left + Inches(0.28), Inches(4.05), Inches(3.4), Inches(2.0), body, 14, False, MUTED)
    footer(s, 3)
    notes(s, "Signpost time: ~4 minutes on the package, ~5 on gaps, ~5 on approach and ask. Invite questions at the end unless a hard-gate concern is raised.")

    # ------------------------------------------------------------------ 04 Monolith
    s = blank()
    content_head(s, "01  ·  Inherited package", "A single offline repository — two layers of authority",
                 "No hidden services. No answer key. All 84 injects disclosed on day one.")
    # two columns
    rect(s, Inches(0.55), Inches(1.95), Inches(6.0), Inches(4.85), WHITE, RULE)
    rect(s, Inches(0.55), Inches(1.95), Inches(6.0), Inches(0.55), NAVY)
    put(s, Inches(0.75), Inches(2.05), Inches(5.6), Inches(0.4), "IMMUTABLE  —  CHALLENGE EVIDENCE", 13, True, WHITE, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    left_items = [
        "case/  ·  integrated narrative and stakeholder pack",
        "data/  ·  84 injects and cited fixtures",
        "knowledge/  ·  policy extracts, including untrusted SOP",
        "source_documents/  ·  protocol and label extracts",
        "evaluation/  ·  15 public fixtures, not answer keys",
        "requirements/  ·  scoring and evidence standard",
        "FILE_HASHES.csv protects this layer",
    ]
    tf = tf_box(s, Inches(0.8), Inches(2.7), Inches(5.5), Inches(3.85))
    for i, item in enumerate(left_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        _run(p, item, 14, False, INK)

    rect(s, Inches(6.8), Inches(1.95), Inches(6.0), Inches(4.85), WHITE, RULE)
    rect(s, Inches(6.8), Inches(1.95), Inches(6.0), Inches(0.55), TEAL)
    put(s, Inches(7.0), Inches(2.05), Inches(5.6), Inches(0.4), "WRITABLE  —  TEAM 3 ONLY", 13, True, WHITE, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    right_items = [
        "submission/src  ·  controls and advisory workflows",
        "submission/tests  ·  deterministic hard-gate suite",
        "submission/evaluation  ·  TEVV harness and graders",
        "submission/artefacts  ·  00–30 working artefacts",
        "submission/evidence  ·  inspectable machine pack",
        "submission/runbooks  ·  setup, ops, incident, AI-off",
        "Nothing here may silently rewrite challenge files",
    ]
    tf = tf_box(s, Inches(7.05), Inches(2.7), Inches(5.5), Inches(3.85))
    for i, item in enumerate(right_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        _run(p, item, 14, False, INK)
    footer(s, 4)
    notes(s, "Emphasise: monolith means one inspectable package, not one conflated authority. If a reviewer cannot tell challenge fact from Team 3 interpretation, we have failed. FILE_HASHES is the boundary.")

    # ------------------------------------------------------------------ 05 Constraints
    s = blank()
    content_head(s, "01  ·  Inherited package", "What the package already forbids — and why that is the product",
                 "AEGIS is an advisory overlay. Expanding write-path would not close a gap; it would create a scoring and GxP failure.")
    items = [
        ("Three workflows only", "A  Batch evidence pack\nB  PV intake / signal support\nC  Supply options / cold-chain"),
        ("Never autonomous", "Disposition, final PV judgement,\nallocation / shipment, quality-status\nchange, recall, eligibility"),
        ("Fail closed", "Stale authorisation, untrusted SOP,\npoisoned tool, silent unit convert,\nfabricated citation"),
        ("Default runtime", "Deterministic, AI-disabled,\noffline fixtures. Inference is\noptional, budgeted, kill-switched"),
    ]
    for i, (title, body) in enumerate(items):
        left = Inches(0.55) + Inches(i * 3.15)
        rect(s, left, Inches(1.95), Inches(3.0), Inches(4.85), WHITE, RULE)
        rect(s, left, Inches(1.95), Inches(3.0), Inches(0.08), GOLD)
        put(s, left + Inches(0.18), Inches(2.25), Inches(2.65), Inches(1.1), title, 16, True, NAVY)
        put(s, left + Inches(0.18), Inches(3.5), Inches(2.65), Inches(2.9), body, 13, False, MUTED)
    footer(s, 5)
    notes(s, "Stakeholders sometimes ask ‘why not just automate eligibility / recall?’ Answer: D-001 / D-203 / INJ-006. The board target does not waive Quality authority. This slide is the product boundary.")

    # ------------------------------------------------------------------ 06 Gap section
    s = blank()
    rect(s, Inches(0), Inches(0), W, H, NAVY)
    rect(s, Inches(0), Inches(0), Inches(0.14), H, GOLD)
    put(s, Inches(0.75), Inches(2.55), Inches(11.5), Inches(0.35), "02", 14, True, GOLD)
    put(s, Inches(0.75), Inches(3.0), Inches(11.5), Inches(0.8), "The gaps that mattered", 36, True, WHITE)
    put(s, Inches(0.75), Inches(4.0), Inches(11.2), Inches(1.1),
        "The catalogue was complete. The participant layer was not yet inspectable.\nThose are different problems.", 18, False, RGBColor(0xC5, 0xCD, 0xD6))
    notes(s, "Pause. The next slides distinguish coverage holes we own from challenge contradictions we must not ‘fix’.")

    # ------------------------------------------------------------------ 07 Taxonomy
    s = blank()
    content_head(s, "02  ·  Gaps", "Four classes of gap — only two were ours to close by building",
                 "Treating all four as ‘missing data’ would have been the wrong programme.")
    rows = [
        ("A", "Coverage", "Close", "Injects with no owning artefact row and no executed test. Invisible if one treats the design map as proof."),
        ("B", "Evidence pack", "Close", "Required files existed, but a reviewer could not see per-test IDs, a current evaluate run, or an A/B/C citation export."),
        ("C", "Case ≠ data", "Record", "Challenge narrative and cited files disagree. Inventing the missing row would be a fabrication hard-gate."),
        ("D", "Residual product", "Label", "Four public fixtures remain not-implemented. Some planned TEST-A/B/C IDs are still incremental. Scorecard prose can lag."),
    ]
    shape = s.shapes.add_table(5, 4, Inches(0.55), Inches(1.95), Inches(12.25), Inches(4.7))
    tbl = shape.table
    tbl.columns[0].width = Inches(0.7)
    tbl.columns[1].width = Inches(2.3)
    tbl.columns[2].width = Inches(1.6)
    tbl.columns[3].width = Inches(7.65)
    headers = ["", "Class", "Stance", "Meaning for this programme"]
    for c, h in enumerate(headers):
        set_cell(tbl, 0, c, h, True)
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            set_cell(tbl, r, c, val)
    style_table(tbl)
    footer(s, 7)
    notes(s, "Spend time on C. Stakeholders will want the data ‘fixed’. Hold the line: PACKAGE_SCOPE says contradictions are conditions. A-501 to A-504 record them.")

    # ------------------------------------------------------------------ 08 Coverage numbers
    s = blank()
    content_head(s, "02  ·  Gaps", "Coverage before close-out — against the 84-inject catalogue",
                 "Source of truth: data/injects.json. EVIDENCE_MAP.md was not used as a coverage source.")
    kpis = [
        ("38", "Implemented", "Named in code or tests", TEAL),
        ("31", "Documented only", "Artefact mention, no test", GOLD),
        ("4", "Bucket only", "A range, not an ID", ROSE),
        ("11", "Uncovered", "No deliverable citation", ROSE),
    ]
    for i, (n, label, sub, color) in enumerate(kpis):
        left = Inches(0.55) + Inches(i * 3.15)
        rect(s, left, Inches(1.95), Inches(3.0), Inches(2.15), WHITE, RULE)
        put(s, left + Inches(0.15), Inches(2.1), Inches(2.7), Inches(0.7), n, 36, True, color, PP_ALIGN.CENTER)
        put(s, left + Inches(0.15), Inches(2.8), Inches(2.7), Inches(0.35), label, 14, True, NAVY, PP_ALIGN.CENTER)
        put(s, left + Inches(0.15), Inches(3.2), Inches(2.7), Inches(0.6), sub, 12, False, MUTED, PP_ALIGN.CENTER)

    rect(s, Inches(0.55), Inches(4.35), Inches(6.0), Inches(2.5), WHITE, RULE)
    put(s, Inches(0.75), Inches(4.5), Inches(5.6), Inches(0.35), "Uncovered — 11", 14, True, ROSE)
    put(s, Inches(0.75), Inches(4.95), Inches(5.6), Inches(1.7),
        "Research  INJ-007, 009–012\nClinical write-path  INJ-015–017, 019, 020\nPrivacy  INJ-062 (skipped in artefact 17)\n\nThese sit outside Workflows A–C write path (D-203).",
        13, False, INK)

    rect(s, Inches(6.8), Inches(4.35), Inches(6.0), Inches(2.5), WHITE, RULE)
    put(s, Inches(7.0), Inches(4.5), Inches(5.6), Inches(0.35), "Bucket only — 4", 14, True, GOLD)
    put(s, Inches(7.0), Inches(4.95), Inches(5.6), Inches(1.7),
        "INJ-027  PAT model vs batch recipe\nINJ-049  Variation classification dispute\nINJ-053  Counterfeit / returns\nINJ-055  CMO capacity over-promise",
        13, False, INK)
    footer(s, 8)
    notes(s, "The 11 were not ‘missing datasets’. Files existed. What was missing was a named control and a test. D-203 already said we must not build discovery or eligibility automation — so the correct close is abstain, not a new workflow.")

    # ------------------------------------------------------------------ 09 Evidence + code
    s = blank()
    content_head(s, "02  ·  Gaps", "Inspectability failures a reviewer would have found in the room",
                 "The four mandatory evidence files were present. Their content would not have survived a close reading.")
    rect(s, Inches(0.55), Inches(1.95), Inches(6.0), Inches(4.85), WHITE, RULE)
    put(s, Inches(0.75), Inches(2.15), Inches(5.6), Inches(0.4), "Evidence pack", 16, True, NAVY)
    tf = tf_box(s, Inches(0.75), Inches(2.7), Inches(5.6), Inches(3.85))
    for i, item in enumerate([
        "test_results.json recorded one ALL row for 60 tests",
        "evaluation_results.json was four days stale",
        "No exported A/B/C citation pack (FR-X-05)",
        "Phase 0–4 exit note still claimed 35 tests",
        "Scorecard artefacts still quoted 46 tests",
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        _run(p, "—  " + item, 14, False, INK)

    rect(s, Inches(6.8), Inches(1.95), Inches(6.0), Inches(4.85), WHITE, RULE)
    put(s, Inches(7.0), Inches(2.15), Inches(5.6), Inches(0.4), "Implementation", 16, True, NAVY)
    tf = tf_box(s, Inches(7.0), Inches(2.7), Inches(5.6), Inches(3.85))
    for i, item in enumerate([
        "Protocol helper invented versions 3.1 / 4.0",
        "Fixtures are global 5.0 and site IN-014 on 4.1",
        "Named automated tests covered 006, 066, 067 only",
        "Dashboard and runbooks carried no real inject IDs",
        "Some evaluation cases used inject IDs as labels",
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        _run(p, "—  " + item, 14, False, INK)
    footer(s, 9)
    notes(s, "This is the ‘we would have failed a cold review’ slide. Stay factual. Then move to what we will not fabricate — that earns trust.")

    # ------------------------------------------------------------------ 10 Tensions
    s = blank()
    content_head(s, "02  ·  Gaps", "Challenge tensions we will not resolve by invention",
                 "81 of 84 cited files instantiate the case sentence. Three do not. Hash drift is recorded, not overwritten.")
    shape = s.shapes.add_table(5, 3, Inches(0.55), Inches(1.95), Inches(12.25), Inches(4.7))
    tbl = shape.table
    tbl.columns[0].width = Inches(1.7)
    tbl.columns[1].width = Inches(1.4)
    tbl.columns[2].width = Inches(9.15)
    for c, h in enumerate(["Inject", "Record", "What a reviewer will see — and what we will not do"]):
        set_cell(tbl, 0, c, h, True)
    data = [
        ("INJ-002", "A-501", "Case names Supply service level. kpi_conflicts.csv has Manufacturing, Quality, Safety, Clinical — no Supply row. We do not invent one."),
        ("INJ-041", "A-502", "sensitive_segments.csv keys pregnancy/minor on PV-1020. icsr_cases.csv has PV-1001, 1009, 1014 only. We do not fabricate ICSR PV-1020."),
        ("INJ-046", "A-503", "Case: US pending and two absent leaflets. Files: EU/US/IN approved and authorised, versions 6/5/3. We cite what is present."),
        ("INJ-065 / 066", "A-504", "On-disk SHA-256 differs from FILE_HASHES.csv. The inject content is still present. We do not overwrite challenge files."),
    ]
    for r, row in enumerate(data, start=1):
        for c, val in enumerate(row):
            set_cell(tbl, r, c, val)
    style_table(tbl)
    footer(s, 10)
    notes(s, "If challenged ‘just add the Supply KPI’, answer: that would be fabricated evidence — a non-waivable hard gate. Honesty is the control.")

    # ------------------------------------------------------------------ 11 Approach divider
    s = blank()
    rect(s, Inches(0), Inches(0), W, H, NAVY)
    rect(s, Inches(0), Inches(0), Inches(0.14), H, GOLD)
    put(s, Inches(0.75), Inches(2.55), Inches(11.5), Inches(0.35), "03", 14, True, GOLD)
    put(s, Inches(0.75), Inches(3.0), Inches(11.5), Inches(0.8), "The close-out approach", 36, True, WHITE)
    put(s, Inches(0.75), Inches(4.0), Inches(11.2), Inches(1.1),
        "Smallest complete path. No new regulated write path.\nNo synthetic datasets.", 18, False, RGBColor(0xC5, 0xCD, 0xD6))
    notes(s, "Transition: method first, then status, then the ask.")

    # ------------------------------------------------------------------ 12 Method
    s = blank()
    content_head(s, "03  ·  Approach", "The bar we applied — then the sequence",
                 "A gap is closed only when a reviewer can find the ID, the requirement, and a passing deterministic test.")
    rect(s, Inches(0.55), Inches(1.95), Inches(12.25), Inches(1.35), WHITE, RULE)
    put(s, Inches(0.8), Inches(2.1), Inches(11.8), Inches(0.3), "COVERAGE BAR", 11, True, TEAL)
    put(s, Inches(0.8), Inches(2.5), Inches(11.8), Inches(0.55),
        "Named in the owning artefact    →    traced to a TEST-ID in artefact 09    →    executed by a deterministic test.\nChallenge files are cited. They are never rewritten.",
        15, False, INK)

    steps = [
        ("1", "Bound", "Research and clinical write-path injects abstain. D-203 stands. No eligibility or discovery automation."),
        ("2", "Name", "Each former hole is written into 05, 07, 09 or 17 — individual IDs, not ranges."),
        ("3", "Execute", "One control module evaluates all 84. Workflows A/C and privacy surface 027, 049, 053, 055, 062."),
        ("4", "Prove", "Per-test results, inject register, A/B/C audit export, current evaluate run, regenerated hashes."),
    ]
    for i, (n, title, body) in enumerate(steps):
        left = Inches(0.55) + Inches(i * 3.15)
        rect(s, left, Inches(3.55), Inches(3.0), Inches(3.25), WHITE, RULE)
        put(s, left + Inches(0.18), Inches(3.7), Inches(2.65), Inches(0.4), n, 22, True, GOLD)
        put(s, left + Inches(0.18), Inches(4.2), Inches(2.65), Inches(0.4), title, 16, True, NAVY)
        put(s, left + Inches(0.18), Inches(4.7), Inches(2.65), Inches(1.85), body, 13, False, MUTED)
    footer(s, 12)
    notes(s, "This is the method they are being asked to endorse. If they want a different bar — e.g. invent missing rows — that is a hard-gate discussion, not a documentation tweak.")

    # ------------------------------------------------------------------ 13 Status
    s = blank()
    content_head(s, "03  ·  Approach", "Current status — closed versus labelled residual",
                 "Live machine evidence, 16 August 2026. Advisory POC. Not a production recommendation.")
    rect(s, Inches(0.55), Inches(1.95), Inches(6.0), Inches(4.85), WHITE, RULE)
    rect(s, Inches(0.55), Inches(1.95), Inches(0.1), Inches(4.85), OK)
    put(s, Inches(0.85), Inches(2.15), Inches(5.5), Inches(0.4), "Closed", 18, True, OK)
    tf = tf_box(s, Inches(0.85), Inches(2.7), Inches(5.45), Inches(3.85))
    for i, item in enumerate([
        "All 84 injects have an executable control",
        "The 15 former holes are named and tested",
        "Protocol versions read from fixtures (5.0 / 4.1)",
        "test_results.json: 60 / 60, one row per test",
        "evaluate.py: 63 pass, 0 fail, 4 labelled NI",
        "Audit export for Workflows A, B and C",
        "Structure gate and submission hashes pass",
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        _run(p, "—  " + item, 14, False, INK)

    rect(s, Inches(6.8), Inches(1.95), Inches(6.0), Inches(4.85), WHITE, RULE)
    rect(s, Inches(6.8), Inches(1.95), Inches(0.1), Inches(4.85), GOLD)
    put(s, Inches(7.1), Inches(2.15), Inches(5.5), Inches(0.4), "Labelled residual", 18, True, GOLD)
    tf = tf_box(s, Inches(7.1), Inches(2.7), Inches(5.45), Inches(3.85))
    for i, item in enumerate([
        "PUB-10, 12, 14, 15 remain not_implemented",
        "POC threshold allows labelled NI; zero hard-gate fails",
        "Some planned TEST-A/B/C IDs still incremental",
        "Artefacts 22 / 28 / 30 may still quote older counts",
        "Dashboard remains thin on named inject IDs",
        "Nine data/*.csv files are unused by any inject",
        "No production go. Conditional POC only (D-405)",
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        _run(p, "—  " + item, 14, False, INK)
    footer(s, 13)
    notes(s, "Be precise on NI: reliability, integration, FinOps, clinical public fixtures. They are indexed and labelled. Pilot would require 15/15. We are not claiming pilot.")

    # ------------------------------------------------------------------ 14 Ask
    s = blank()
    content_head(s, "Decision requested", "What we need from this room")
    rect(s, Inches(0.55), Inches(1.85), Inches(12.25), Inches(1.55), WHITE, RULE)
    put(s, Inches(0.8), Inches(2.05), Inches(11.8), Inches(0.35), "RECOMMENDATION", 11, True, TEAL)
    put(s, Inches(0.8), Inches(2.45), Inches(11.8), Inches(0.7),
        "Accept the close-out as conditional-go for the advisory POC. Do not authorise production. Do not authorise invention of missing challenge rows.",
        16, False, INK)

    asks = [
        ("Endorse the bar", "Named artefact + TEST-ID + deterministic test is the definition of ‘covered’."),
        ("Accept recorded tensions", "A-501 to A-504 stand. Challenge contradictions remain visible."),
        ("Accept labelled NI", "PUB-10 / 12 / 14 / 15 stay not_implemented until a later increment."),
        ("Hold the write-path line", "No eligibility, disposition, final PV, allocate, ship, or recall — even under schedule pressure."),
    ]
    for i, (title, body) in enumerate(asks):
        left = Inches(0.55) + Inches((i % 2) * 6.35)
        top = Inches(3.65) + Inches((i // 2) * 1.6)
        rect(s, left, top, Inches(6.15), Inches(1.45), WHITE, RULE)
        put(s, left + Inches(0.25), top + Inches(0.18), Inches(5.7), Inches(0.35), title, 15, True, NAVY)
        put(s, left + Inches(0.25), top + Inches(0.6), Inches(5.7), Inches(0.65), body, 13, False, MUTED)
    footer(s, 14)
    notes(s, "Close on the four asks. If time remains, offer to walk INJECT_CONTROL_REGISTER.md or audit_export.json live. Do not reopen EVIDENCE_MAP as the coverage source.")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    return OUT


if __name__ == "__main__":
    print(f"wrote {build()}")
