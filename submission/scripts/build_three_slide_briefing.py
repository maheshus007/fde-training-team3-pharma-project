#!/usr/bin/env python3
"""Three-slide stakeholder briefing: monolith, coupling gaps, multi-agent close-out."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parents[1] / "artefacts" / "AEGIS_THREE_SLIDE_BRIEFING.pptx"

NAVY = RGBColor(0x0A, 0x16, 0x28)
TEAL = RGBColor(0x0E, 0x4D, 0x5A)
GOLD = RGBColor(0xB8, 0x95, 0x2C)
PAPER = RGBColor(0xF4, 0xF1, 0xEA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1A, 0x1F, 0x26)
MUTED = RGBColor(0x5A, 0x62, 0x70)
RULE = RGBColor(0xD8, 0xD1, 0xC4)
ROSE = RGBColor(0x7A, 0x2E, 0x2E)
OK = RGBColor(0x1F, 0x5C, 0x45)

W = Inches(13.333)
H = Inches(7.5)


def _run(p, text, size, bold=False, color=INK):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return r


def box(slide, l, t, w, h):
    sh = slide.shapes.add_textbox(l, t, w, h)
    tf = sh.text_frame
    tf.word_wrap = True
    return tf


def put(slide, l, t, w, h, text, size=16, bold=False, color=INK, align=PP_ALIGN.LEFT):
    tf = box(slide, l, t, w, h)
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    _run(p, text, size, bold, color)
    return tf


def rect(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    else:
        sh.line.fill.background()
    return sh


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def footer(slide, n):
    rect(slide, Inches(0), Inches(7.22), W, Inches(0.28), NAVY)
    put(slide, Inches(0.5), Inches(7.24), Inches(9), Inches(0.24),
        "AEGIS-PHARMA  ·  Stakeholder briefing  ·  Confidential  ·  16 August 2026",
        10, False, RGBColor(0xC9, 0xD0, 0xD8))
    put(slide, Inches(11.5), Inches(7.24), Inches(1.3), Inches(0.24),
        f"{n}  /  3", 10, False, GOLD, PP_ALIGN.RIGHT)


def head(slide, kicker, title):
    rect(slide, Inches(0), Inches(0), W, Inches(1.18), NAVY)
    rect(slide, Inches(0), Inches(0), Inches(0.12), Inches(1.18), GOLD)
    put(slide, Inches(0.5), Inches(0.16), Inches(12.3), Inches(0.28), kicker.upper(), 11, True, GOLD)
    put(slide, Inches(0.5), Inches(0.48), Inches(12.3), Inches(0.55), title, 24, True, WHITE)


def slide_bg(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, Inches(0), Inches(0), W, H, PAPER)
    return s


def build() -> Path:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # ================================================================== 1
    s = slide_bg(prs)
    head(s, "01  ·  Base repository  ·  monolithic application",
         "One package. One runtime. Three advisory workflows.")
    put(s, Inches(0.5), Inches(1.32), Inches(12.3), Inches(0.36),
        "AEGIS sits on a single offline tree. Challenge evidence is immutable. Team 3 work lives only under submission/.",
        13, False, MUTED)

    # How it works - flow
    steps = [
        ("1", "Request", "User, purpose,\nobject, role"),
        ("2", "Bind", "Execution-time\nauthz  ·  deny stale"),
        ("3", "Workflow", "A  Batch   B  PV\nC  Supply options"),
        ("4", "Guard", "policy_guard\ncontracts  ·  fail closed"),
        ("5", "Pack", "Citations, conflicts,\nabstentions  ·  no write"),
    ]
    for i, (n, title, body) in enumerate(steps):
        left = Inches(0.5) + Inches(i * 2.52)
        rect(s, left, Inches(1.8), Inches(2.38), Inches(2.15), WHITE, RULE)
        put(s, left + Inches(0.12), Inches(1.9), Inches(0.4), Inches(0.35), n, 18, True, GOLD)
        put(s, left + Inches(0.5), Inches(1.92), Inches(1.75), Inches(0.35), title, 14, True, NAVY)
        put(s, left + Inches(0.12), Inches(2.4), Inches(2.14), Inches(1.4), body, 12, False, MUTED)
        if i < 4:
            put(s, left + Inches(2.22), Inches(2.55), Inches(0.3), Inches(0.35), "→", 16, True, GOLD)

    # Three workflow cards
    wfs = [
        ("Workflow A  ·  Batch", "Reconcile MES, LIMS, warehouse and QP packet. Surface genealogy, OOS/OOT and unit conflicts. Readiness only — no disposition."),
        ("Workflow B  ·  PV", "Intake, duplicate cluster, clocks, listedness by source. Advisory signal support — no final causality or reportability."),
        ("Workflow C  ·  Supply", "Shortage and cold-chain options under MA and hold constraints. Draft options — no reserve, allocate, ship or recall."),
    ]
    for i, (title, body) in enumerate(wfs):
        left = Inches(0.5) + Inches(i * 4.2)
        rect(s, left, Inches(4.15), Inches(4.05), Inches(2.7), WHITE, RULE)
        rect(s, left, Inches(4.15), Inches(4.05), Inches(0.07), TEAL)
        put(s, left + Inches(0.2), Inches(4.28), Inches(3.7), Inches(0.4), title, 14, True, NAVY)
        put(s, left + Inches(0.2), Inches(4.75), Inches(3.7), Inches(1.9), body, 12, False, INK)

    footer(s, 1)
    notes(s, "Open with how the monolith works today: one runtime, three advisory workflows, shared policy_guard. Everything shares one code path and one evidence tree. That is the strength for GxP inspectability — and the coupling problem on the next slide.")

    # ================================================================== 2
    s = slide_bg(prs)
    head(s, "02  ·  Gaps in the repository",
         "Tightly coupled — and 84 injects were never first-class.")
    put(s, Inches(0.5), Inches(1.32), Inches(12.3), Inches(0.36),
        "The catalogue discloses INJ-001…084 on day one. Implementation was bolted onto A, B and C. Injects with no home were skipped or mentioned as a range.",
        13, False, MUTED)

    # Coupling problem
    rect(s, Inches(0.5), Inches(1.8), Inches(6.15), Inches(5.05), WHITE, RULE)
    put(s, Inches(0.7), Inches(1.95), Inches(5.8), Inches(0.35), "Tight coupling", 16, True, ROSE)
    tf = box(s, Inches(0.7), Inches(2.4), Inches(5.8), Inches(4.2))
    items = [
        "One orchestrator owns batch, PV and supply.",
        "Shared modules decide every inject in the same path.",
        "Research and clinical facts have no workflow of their own.",
        "A change in batch logic can shift PV or supply behaviour.",
        "Dashboard and runbooks inherited the same coupling — few named inject IDs.",
        "Result: the platform cannot grow a domain without touching the others.",
    ]
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        _run(p, "—  " + item, 13, False, INK)

    rect(s, Inches(6.85), Inches(1.8), Inches(6.0), Inches(5.05), WHITE, RULE)
    put(s, Inches(7.05), Inches(1.95), Inches(5.6), Inches(0.35), "84 injects — not implemented as workflows", 15, True, ROSE)

    kpis = [("38", "In A / B / C code"), ("31", "Documented only"), ("4", "Range only"), ("11", "No home at all")]
    for i, (n, lab) in enumerate(kpis):
        top = Inches(2.45) + Inches(i * 0.72)
        put(s, Inches(7.1), top, Inches(1.1), Inches(0.55), n, 22, True, ROSE if i >= 2 else GOLD)
        put(s, Inches(8.3), top + Inches(0.08), Inches(4.3), Inches(0.5), lab, 14, False, INK)

    put(s, Inches(7.1), Inches(5.4), Inches(5.5), Inches(1.2),
        "The 11 with no home are Discovery (assay, omics, models, target evidence) and Clinical (IRT, unblinding, consent, adjudication, site risk). They cannot be forced into batch, PV or supply without coupling them further.",
        12, False, MUTED)

    footer(s, 2)
    notes(s, "The gap is structural: tight coupling plus incomplete inject implementation. Do not say the data was missing — the files exist. Say the monolith had only three workflow slots, so 84 injects could not all be implemented as independent capabilities.")

    # ================================================================== 3
    s = slide_bg(prs)
    head(s, "03  ·  How we cover the gaps",
         "Multi-agent platform — workflows defined independently.")
    put(s, Inches(0.5), Inches(1.32), Inches(12.3), Inches(0.36),
        "Each workflow is its own agent. A shared platform holds policy, contracts, entitlements and evidence. Uncovered injects get two new agents: Discovery and Clinical.",
        13, False, MUTED)

    # Platform bar
    rect(s, Inches(0.5), Inches(1.78), Inches(12.35), Inches(0.7), NAVY)
    put(s, Inches(0.7), Inches(1.9), Inches(12.0), Inches(0.45),
        "Shared platform    ·    policy_guard    ·    versioned contracts    ·    execution-time authz    ·    evidence register    ·    AI-disabled continuity",
        13, True, WHITE, PP_ALIGN.CENTER)

    agents = [
        ("A  Batch", "Existing", TEAL, "Genealogy, OOS, units, QP gaps"),
        ("B  PV", "Existing", TEAL, "Intake, clocks, listedness, segments"),
        ("C  Supply", "Existing", TEAL, "Cold-chain, shortage, CMO, returns"),
        ("D  Discovery", "New", GOLD, "Assay, omics, models, target conflict\nINJ-007, 009–012"),
        ("E  Clinical", "New", GOLD, "Protocol, IRT, consent, endpoints\nINJ-013–017, 019–020"),
    ]
    for i, (title, tag, color, body) in enumerate(agents):
        left = Inches(0.5) + Inches(i * 2.52)
        rect(s, left, Inches(2.65), Inches(2.38), Inches(2.55), WHITE, RULE)
        rect(s, left, Inches(2.65), Inches(2.38), Inches(0.07), color)
        put(s, left + Inches(0.12), Inches(2.78), Inches(2.14), Inches(0.4), title, 13, True, NAVY)
        put(s, left + Inches(0.12), Inches(3.18), Inches(2.14), Inches(0.28), tag.upper(), 10, True, color)
        put(s, left + Inches(0.12), Inches(3.5), Inches(2.14), Inches(1.5), body, 11, False, MUTED)

    put(s, Inches(0.5), Inches(5.35), Inches(12.3), Inches(0.28),
        "Independence rule", 12, True, TEAL)
    put(s, Inches(0.5), Inches(5.65), Inches(12.3), Inches(1.35),
        "Agents do not share a single if/else. They publish citation packs on the same contract. Discovery surfaces comparability and model-trust conflicts — it does not promote a research model. Clinical surfaces protocol, consent and clock conflicts — it does not decide eligibility, unblinding or endpoint close. Hard gates stay on the platform, not inside any one agent.",
        13, False, INK)

    footer(s, 3)
    notes(s, "Close: we are not splitting the repo into five products. We are splitting workflow ownership. Discovery and Clinical are the two missing agents for the 11 uncovered injects. They remain advisory. Ask stakeholders to endorse independent workflow definition on a shared multi-agent platform — not a larger monolith.")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    return OUT


if __name__ == "__main__":
    print(f"wrote {build()}")
